from __future__ import annotations

from pathlib import Path
from uuid import uuid4

from pypdf import PdfReader
from pptx import Presentation

from app.db import SessionLocal
from app.models import AssetStatus, AssetType, ContentChunk, SourceAsset
from app.services.embeddings import embed_text
from app.services.generation import generate_course_materials
from app.services.text_utils import clean_text, split_for_chunks


def _extract_pdf_sections(path: Path) -> list[tuple[str, str]]:
    reader = PdfReader(str(path))
    sections: list[tuple[str, str]] = []

    for page_number, page in enumerate(reader.pages, start=1):
        text = clean_text(page.extract_text() or "")
        for index, chunk in enumerate(split_for_chunks(text), start=1):
            label = f"Page {page_number}" if index == 1 else f"Page {page_number} · Part {index}"
            sections.append((label, chunk))

    if not sections:
        raise ValueError("No readable text was found in the PDF.")

    return sections


def _extract_slide_text(presentation: Presentation, slide_index: int) -> str:
    slide = presentation.slides[slide_index]
    text_bits: list[str] = []

    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text_bits.append(shape.text)

    notes_text = ""
    if slide.has_notes_slide:
        for shape in slide.notes_slide.shapes:
            if hasattr(shape, "text") and shape.text:
                notes_text += "\n" + shape.text

    return clean_text("\n".join(text_bits) + "\n" + notes_text)


def _extract_pptx_sections(path: Path) -> list[tuple[str, str]]:
    presentation = Presentation(str(path))
    sections: list[tuple[str, str]] = []

    for slide_number in range(len(presentation.slides)):
        text = _extract_slide_text(presentation, slide_number)
        if not text:
            continue
        for index, chunk in enumerate(split_for_chunks(text), start=1):
            label = f"Slide {slide_number + 1}" if index == 1 else f"Slide {slide_number + 1} · Part {index}"
            sections.append((label, chunk))

    if not sections:
        raise ValueError("No readable text was found in the PowerPoint.")

    return sections


def extract_sections(asset: SourceAsset) -> list[tuple[str, str]]:
    path = Path(asset.storage_path)
    if asset.type == AssetType.PDF:
        return _extract_pdf_sections(path)
    if asset.type == AssetType.PPTX:
        return _extract_pptx_sections(path)
    raise ValueError(f"Unsupported asset type: {asset.type.value}")


def process_asset(asset_id: str) -> None:
    db = SessionLocal()
    try:
        asset = db.get(SourceAsset, asset_id)
        if asset is None:
            return

        asset.status = AssetStatus.PROCESSING
        db.commit()

        sections = extract_sections(asset)

        existing_chunks = db.query(ContentChunk).filter(ContentChunk.asset_id == asset.id).all()
        for chunk in existing_chunks:
            db.delete(chunk)
        db.flush()

        for index, (source_ref, text) in enumerate(sections, start=1):
            db.add(
                ContentChunk(
                    id=str(uuid4()),
                    course_id=asset.course_id,
                    asset_id=asset.id,
                    sequence_index=index,
                    source_ref=source_ref,
                    text=text,
                    embedding=embed_text(text),
                    topic_ids=[],
                )
            )

        db.flush()
        asset.status = AssetStatus.READY
        asset.details = {
            "section_count": len(sections),
            "source_type": asset.type.value,
            "last_processed_sections": [source_ref for source_ref, _text in sections[:5]],
        }

        generate_course_materials(db, asset.course_id)
        db.commit()
    except Exception as exc:  # noqa: BLE001
        asset = db.get(SourceAsset, asset_id)
        if asset is not None:
            asset.status = AssetStatus.FAILED
            asset.details = {"error": str(exc)}
            db.commit()
    finally:
        db.close()
