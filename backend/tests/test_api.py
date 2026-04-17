from __future__ import annotations

from io import BytesIO

from fastapi.testclient import TestClient
from pptx import Presentation
from reportlab.pdfgen import canvas

from app.db import Base, engine
from app.main import app


client = TestClient(app)
USER_HEADERS = {"x-user-id": "integration-user", "x-user-name": "Integration User"}


def setup_function() -> None:
    Base.metadata.drop_all(bind=engine)
    Base.metadata.create_all(bind=engine)


def _create_pdf() -> bytes:
    stream = BytesIO()
    pdf = canvas.Canvas(stream)
    pdf.drawString(72, 720, "Chemical Equilibrium")
    pdf.drawString(72, 700, "Kp measures the ratio of gaseous products to reactants at equilibrium.")
    pdf.drawString(72, 680, "A higher Kp means products are favored.")
    pdf.save()
    return stream.getvalue()


def _create_pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "RSA Encryption"
    slide.placeholders[1].text = "Public keys encrypt data. Private keys decrypt it."
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _create_overview_pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Block 4 Overview"
    slide.placeholders[1].text = "Kp describes equilibrium for gases. Larger Kp values favor products at equilibrium."
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def test_pdf_upload_generates_study_materials() -> None:
    create_course = client.post(
        "/api/courses",
        json={"title": "Chemistry Midterm", "subject": "Chemistry"},
        headers=USER_HEADERS,
    )
    course_id = create_course.json()["id"]

    upload = client.post(
        f"/api/courses/{course_id}/assets",
        files={"file": ("equilibrium.pdf", _create_pdf(), "application/pdf")},
        headers=USER_HEADERS,
    )
    assert upload.status_code == 201

    status_response = client.get(f"/api/courses/{course_id}/ingestion-status", headers=USER_HEADERS)
    assert status_response.status_code == 200
    assert status_response.json()["ready_assets"] == 1

    overview = client.get(f"/api/courses/{course_id}/overview", headers=USER_HEADERS)
    payload = overview.json()
    assert payload["learning_counts"]["flashcard"] >= 1
    assert payload["topics"][0]["source_refs"][0]["location"].startswith("Page")


def test_pptx_upload_creates_slide_references() -> None:
    create_course = client.post(
        "/api/courses",
        json={"title": "Security", "subject": "Computer Science"},
        headers=USER_HEADERS,
    )
    course_id = create_course.json()["id"]

    upload = client.post(
        f"/api/courses/{course_id}/assets",
        files={"file": ("rsa-intro.pptx", _create_pptx(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        headers=USER_HEADERS,
    )
    assert upload.status_code == 201

    flashcards = client.get(f"/api/courses/{course_id}/flashcards", headers=USER_HEADERS)
    assert flashcards.status_code == 200
    assert "Slide 1" in flashcards.json()[0]["source_refs"][0]["location"]


def test_attempts_update_guided_learning_priority() -> None:
    create_course = client.post(
        "/api/courses",
        json={"title": "Thermo", "subject": "Chemistry"},
        headers=USER_HEADERS,
    )
    course_id = create_course.json()["id"]

    client.post(
        f"/api/courses/{course_id}/assets",
        files={"file": ("equilibrium.pdf", _create_pdf(), "application/pdf")},
        headers=USER_HEADERS,
    )

    quiz_session = client.post(
        f"/api/courses/{course_id}/quiz/session",
        json={"item_count": 1},
        headers=USER_HEADERS,
    )
    item = quiz_session.json()["items"][0]

    attempt = client.post(
        f"/api/learning-items/{item['id']}/attempt",
        json={"response": "I am not sure", "confidence_1_5": 1},
        headers=USER_HEADERS,
    )
    assert attempt.status_code == 200
    assert attempt.json()["priority_score"] >= 70

    guided = client.get(f"/api/courses/{course_id}/guided-learning", headers=USER_HEADERS)
    assert guided.status_code == 200
    assert guided.json()["items"][0]["dynamic_priority"] >= 70


def test_guided_learning_uses_content_not_block_labels() -> None:
    create_course = client.post(
        "/api/courses",
        json={"title": "Equilibrium", "subject": "Chemistry"},
        headers=USER_HEADERS,
    )
    course_id = create_course.json()["id"]

    upload = client.post(
        f"/api/courses/{course_id}/assets",
        files={"file": ("block-4-overview.pptx", _create_overview_pptx(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        headers=USER_HEADERS,
    )
    assert upload.status_code == 201

    guided = client.get(f"/api/courses/{course_id}/guided-learning", headers=USER_HEADERS)
    assert guided.status_code == 200
    item = guided.json()["items"][0]
    assert "block" not in item["learning_item"]["prompt"].lower()
    assert "matter in this course" not in item["learning_item"]["prompt"].lower()
    assert "kp" in item["learning_item"]["prompt"].lower()


def test_guided_learning_attempt_returns_feedback() -> None:
    create_course = client.post(
        "/api/courses",
        json={"title": "Equilibrium", "subject": "Chemistry"},
        headers=USER_HEADERS,
    )
    course_id = create_course.json()["id"]

    client.post(
        f"/api/courses/{course_id}/assets",
        files={"file": ("equilibrium.pdf", _create_pdf(), "application/pdf")},
        headers=USER_HEADERS,
    )

    guided = client.get(f"/api/courses/{course_id}/guided-learning", headers=USER_HEADERS)
    item = guided.json()["items"][0]["learning_item"]

    attempt = client.post(
        f"/api/learning-items/{item['id']}/attempt",
        json={"response": "Kp compares products and reactants at equilibrium.", "confidence_1_5": 3},
        headers=USER_HEADERS,
    )
    assert attempt.status_code == 200
    assert "Quick check" in attempt.json()["feedback"]
